"""
PNU-Pathfinder 해커톤 예선 발표 PPT 생성기
실행: python3 scripts/generate_ppt.py
출력: output/발표자료_PNU-Pathfinder.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# ── 색상 팔레트 ──────────────────────────────────────────
PNU_BLUE   = RGBColor(0x00, 0x38, 0x76)   # 부산대 파랑
ACCENT     = RGBColor(0x10, 0xB9, 0x81)   # 민트/에메랄드
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
DARK       = RGBColor(0x0F, 0x17, 0x2A)
LIGHT_GRAY = RGBColor(0xF1, 0xF5, 0xF9)
SLATE      = RGBColor(0x47, 0x55, 0x69)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]


def rgb(r, g, b):
    return RGBColor(r, g, b)


def add_rect(slide, x, y, w, h, fill_color, alpha=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    return shape


def add_text(slide, text, x, y, w, h, font_size=24, bold=False,
             color=WHITE, align=PP_ALIGN.LEFT, font_name="맑은 고딕"):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox


def add_bullet_box(slide, items, x, y, w, h, font_size=20,
                   bullet="•", color=DARK, bg=None, padding=Inches(0.2)):
    if bg:
        add_rect(slide, x, y, w, h, bg)
    txBox = slide.shapes.add_textbox(x + padding, y + padding,
                                      w - padding*2, h - padding*2)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = f"{bullet}  {item}"
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.name = "맑은 고딕"


# ════════════════════════════════════════════════════════
# SLIDE 1 — 표지
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, H, DARK)
add_rect(slide, 0, Inches(2.8), W, Inches(0.06), ACCENT)

add_text(slide, "제7회 PNU 창의융합AI해커톤",
         Inches(1), Inches(1.2), Inches(11), Inches(0.7),
         font_size=22, color=ACCENT, align=PP_ALIGN.CENTER)

add_text(slide, "PNU-Pathfinder",
         Inches(1), Inches(2.0), Inches(11), Inches(1.2),
         font_size=54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(slide, "부산대학교 학생 맞춤형 학업·진로 자기경영 AI 플랫폼",
         Inches(1), Inches(3.1), Inches(11), Inches(0.7),
         font_size=24, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

add_text(slide, "융합트랙  |  자유주제  |  2026",
         Inches(1), Inches(5.8), Inches(11), Inches(0.5),
         font_size=18, color=SLATE, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════
# SLIDE 2 — 문제 정의
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, Inches(1.4), PNU_BLUE)
add_text(slide, "문제 정의", Inches(0.6), Inches(0.35), Inches(10), Inches(0.7),
         font_size=36, bold=True, color=WHITE)
add_text(slide, "PROBLEM", Inches(0.6), Inches(0.9), Inches(10), Inches(0.4),
         font_size=16, color=ACCENT)

add_rect(slide, 0, Inches(1.4), W, H - Inches(1.4), LIGHT_GRAY)

problems = [
    ("학사 시스템", "성적·수강이력·이수학점"),
    ("학과 홈페이지", "교육과정·이수체계도"),
    ("비교과 시스템", "프로그램·마일리지"),
    ("취업지원 사이트", "채용·진로 프로그램"),
    ("외부 사이트", "자격증·어학 일정"),
]

add_text(slide, "부산대 학생이 확인해야 하는 정보 출처",
         Inches(0.5), Inches(1.6), Inches(12), Inches(0.5),
         font_size=22, bold=True, color=DARK)

for i, (title, sub) in enumerate(problems):
    x = Inches(0.5 + i * 2.5)
    add_rect(slide, x, Inches(2.3), Inches(2.2), Inches(1.5), PNU_BLUE)
    add_text(slide, title, x, Inches(2.4), Inches(2.2), Inches(0.5),
             font_size=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, sub, x, Inches(2.95), Inches(2.2), Inches(0.6),
             font_size=14, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

add_rect(slide, Inches(0.5), Inches(4.1), Inches(12.3), Inches(0.06), ACCENT)

add_text(slide, "결과: 학생이 스스로 모든 정보를 탐색·조합·해석해야 함",
         Inches(0.5), Inches(4.3), Inches(12), Inches(0.5),
         font_size=20, bold=True, color=PNU_BLUE)

pain_points = [
    "1·2학년: 무엇을 먼저 들어야 할지 모른다",
    "3·4학년: 졸업요건 누락을 뒤늦게 발견한다",
    "복수전공·편입·복학생: 교육과정 변화로 더 큰 혼란을 겪는다",
    "모든 학생: 진로 목표에 맞는 실행 계획을 스스로 세워야 한다",
]
add_bullet_box(slide, pain_points,
               Inches(0.5), Inches(4.9), Inches(12.3), Inches(2.2),
               font_size=18, color=DARK, bg=WHITE)


# ════════════════════════════════════════════════════════
# SLIDE 3 — 솔루션 개요
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, Inches(1.4), PNU_BLUE)
add_text(slide, "솔루션", Inches(0.6), Inches(0.35), Inches(10), Inches(0.7),
         font_size=36, bold=True, color=WHITE)
add_text(slide, "SOLUTION", Inches(0.6), Inches(0.9), Inches(10), Inches(0.4),
         font_size=16, color=ACCENT)

add_rect(slide, 0, Inches(1.4), W, H - Inches(1.4), DARK)

add_text(slide, "PNU-Pathfinder",
         Inches(0.5), Inches(1.7), Inches(12), Inches(0.9),
         font_size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(slide,
         "분산된 학업·졸업·진로 정보를 통합하고,\n"
         "학생 개인 상황에 맞는 실행 계획을 AI가 제안하는 웹 서비스",
         Inches(1), Inches(2.7), Inches(11), Inches(1.0),
         font_size=22, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

add_rect(slide, Inches(0.5), Inches(3.9), Inches(12.3), Inches(0.06), ACCENT)

features = ["Smart Planner", "Interactive Roadmap",
            "AI 챗봇", "What-if 분석", "이력서 생성"]
for i, feat in enumerate(features):
    x = Inches(0.5 + i * 2.46)
    add_rect(slide, x, Inches(4.2), Inches(2.3), Inches(1.0), ACCENT)
    add_text(slide, feat, x, Inches(4.4), Inches(2.3), Inches(0.6),
             font_size=16, bold=True, color=DARK, align=PP_ALIGN.CENTER)

add_text(slide,
         "규칙 기반 추천 엔진 + 생성형 AI 설명 → 신뢰성 + 설명 가능성 동시 확보",
         Inches(0.5), Inches(5.6), Inches(12.3), Inches(0.5),
         font_size=18, color=ACCENT, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════
# SLIDE 4 — 핵심 기능 (Smart Planner)
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, Inches(1.4), PNU_BLUE)
add_text(slide, "핵심 기능 ①  Smart Planner", Inches(0.6), Inches(0.3),
         Inches(11), Inches(0.7), font_size=32, bold=True, color=WHITE)
add_text(slide, "지능형 수강 시뮬레이터", Inches(0.6), Inches(0.95),
         Inches(11), Inches(0.4), font_size=16, color=ACCENT)

add_rect(slide, 0, Inches(1.4), W, H - Inches(1.4), LIGHT_GRAY)

cards = [
    ("졸업요건 충족률", "영역별 잔여 학점 자동 계산\n프로그레스 바 시각화"),
    ("수강 추천", "선수과목·졸업필수 반영\n다음 학기 최적 과목 제안"),
    ("재수강 우선순위", "C+ 이하 과목 중\n성적 상승 효과 큰 과목 추천"),
    ("학기 시뮬레이션", "졸업까지 몇 학기 남았는지\n시나리오별 예측"),
]
for i, (title, body) in enumerate(cards):
    col = i % 2
    row = i // 2
    x = Inches(0.5 + col * 6.4)
    y = Inches(1.6 + row * 2.4)
    add_rect(slide, x, y, Inches(6.0), Inches(2.1), WHITE)
    add_rect(slide, x, y, Inches(0.12), Inches(2.1), ACCENT)
    add_text(slide, title, x + Inches(0.25), y + Inches(0.1),
             Inches(5.5), Inches(0.5), font_size=20, bold=True, color=PNU_BLUE)
    add_text(slide, body, x + Inches(0.25), y + Inches(0.65),
             Inches(5.5), Inches(1.2), font_size=16, color=SLATE)


# ════════════════════════════════════════════════════════
# SLIDE 5 — 핵심 기능 (AI 챗봇 + What-if)
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, Inches(1.4), PNU_BLUE)
add_text(slide, "핵심 기능 ②  AI 챗봇 & What-if 분석", Inches(0.6), Inches(0.3),
         Inches(11), Inches(0.7), font_size=32, bold=True, color=WHITE)
add_text(slide, "생성형 AI 연동 기능", Inches(0.6), Inches(0.95),
         Inches(11), Inches(0.4), font_size=16, color=ACCENT)

add_rect(slide, 0, Inches(1.4), W, H - Inches(1.4), DARK)

# AI 챗봇
add_rect(slide, Inches(0.4), Inches(1.6), Inches(6.0), Inches(5.5), PNU_BLUE)
add_text(slide, "AI 챗봇", Inches(0.6), Inches(1.7), Inches(5.5), Inches(0.6),
         font_size=24, bold=True, color=ACCENT)
chat_items = [
    '"다음 학기 뭐 들어야 해?" → 개인 이수현황 기반 답변',
    '"졸업 가능해?" → 충족률 + 부족 요건 설명',
    '"재수강 뭐 해야 해?" → 우선순위 추천 + 이유 설명',
    "OpenAI API 연동 — 자연어로 설명",
    "규칙 기반 로직 + AI 설명 분리 구조",
]
add_bullet_box(slide, chat_items, Inches(0.5), Inches(2.4),
               Inches(5.7), Inches(4.2), font_size=16,
               color=WHITE, bg=None)

# What-if
add_rect(slide, Inches(6.9), Inches(1.6), Inches(6.0), Inches(5.5), rgb(0x1E, 0x29, 0x3B))
add_text(slide, "What-if 시나리오", Inches(7.1), Inches(1.7),
         Inches(5.5), Inches(0.6), font_size=24, bold=True, color=ACCENT)
whatif_items = [
    '"경영학과 복수전공 시 몇 학기 더 필요해?"',
    '"전과 시 기존 이수 과목 인정 범위는?"',
    '"복수전공 시 다음 학기 수강 계획은?"',
    "저학년 진로·전공 의사결정 지원",
]
add_bullet_box(slide, whatif_items, Inches(7.0), Inches(2.4),
               Inches(5.7), Inches(4.2), font_size=16,
               color=WHITE, bg=None)


# ════════════════════════════════════════════════════════
# SLIDE 6 — AI 활용 계획
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, Inches(1.4), PNU_BLUE)
add_text(slide, "AI 활용 계획", Inches(0.6), Inches(0.35),
         Inches(11), Inches(0.7), font_size=36, bold=True, color=WHITE)
add_text(slide, "AI UTILIZATION PLAN", Inches(0.6), Inches(0.95),
         Inches(11), Inches(0.4), font_size=16, color=ACCENT)

add_rect(slide, 0, Inches(1.4), W, H - Inches(1.4), LIGHT_GRAY)

# 서비스 내 AI
add_rect(slide, Inches(0.4), Inches(1.6), Inches(6.0), Inches(5.5), WHITE)
add_rect(slide, Inches(0.4), Inches(1.6), Inches(6.0), Inches(0.55), PNU_BLUE)
add_text(slide, "서비스 내 AI 활용", Inches(0.5), Inches(1.65),
         Inches(5.7), Inches(0.4), font_size=18, bold=True, color=WHITE)
svc_items = [
    "맞춤형 수강·재수강 추천 생성",
    "졸업요건·이수체계 자연어 질의응답",
    "추천 결과 설명 가능한 언어로 해설",
    "이력서·자기소개서 초안 자동 생성",
    "커뮤니티 팁 요약 및 인사이트 제공",
]
add_bullet_box(slide, svc_items, Inches(0.5), Inches(2.3),
               Inches(5.7), Inches(4.5), font_size=17, color=DARK)

# 개발 과정 AI
add_rect(slide, Inches(6.9), Inches(1.6), Inches(6.0), Inches(5.5), WHITE)
add_rect(slide, Inches(6.9), Inches(1.6), Inches(6.0), Inches(0.55), ACCENT)
add_text(slide, "개발 과정 AI 활용", Inches(7.0), Inches(1.65),
         Inches(5.7), Inches(0.4), font_size=18, bold=True, color=DARK)
dev_items = [
    "Claude — 문제 정의, 기획, 문서 작성",
    "ChatGPT / Gemini — 아이디어 검토",
    "Cursor — 코드 생성·디버깅·리팩토링",
    "GitHub Copilot — 개발 보조",
    "Midjourney — UI 시안 생성",
]
add_bullet_box(slide, dev_items, Inches(7.0), Inches(2.3),
               Inches(5.7), Inches(4.5), font_size=17, color=DARK)


# ════════════════════════════════════════════════════════
# SLIDE 7 — AI Agent 협업 아키텍처 (핵심 차별화)
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, Inches(1.4), PNU_BLUE)
add_text(slide, "AI Agent 협업 개발 방식", Inches(0.6), Inches(0.3),
         Inches(11), Inches(0.7), font_size=32, bold=True, color=WHITE)
add_text(slide, "Claude Multi-Agent  ×  MCP (Model Context Protocol)",
         Inches(0.6), Inches(0.95), Inches(11), Inches(0.4),
         font_size=16, color=ACCENT)

add_rect(slide, 0, Inches(1.4), W, H - Inches(1.4), DARK)

# PM Agent 박스 (중앙 상단)
add_rect(slide, Inches(4.5), Inches(1.6), Inches(4.3), Inches(1.1), PNU_BLUE)
add_text(slide, "PM Agent", Inches(4.5), Inches(1.65), Inches(4.3), Inches(0.45),
         font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, "기획 · 문서 · 코드리뷰\nMCP: filesystem  github",
         Inches(4.5), Inches(2.1), Inches(4.3), Inches(0.5),
         font_size=12, color=ACCENT, align=PP_ALIGN.CENTER)

# 화살표 텍스트
add_text(slide, "태스크 분배  ↙                    ↘",
         Inches(2.5), Inches(2.85), Inches(8.5), Inches(0.4),
         font_size=16, color=SLATE, align=PP_ALIGN.CENTER)

# Backend Agent 박스
add_rect(slide, Inches(0.4), Inches(3.4), Inches(5.6), Inches(2.8), rgb(0x1E, 0x29, 0x3B))
add_rect(slide, Inches(0.4), Inches(3.4), Inches(5.6), Inches(0.5), ACCENT)
add_text(slide, "Backend Agent  (Claude Sonnet 4.6)",
         Inches(0.5), Inches(3.43), Inches(5.4), Inches(0.42),
         font_size=15, bold=True, color=DARK, align=PP_ALIGN.CENTER)

be_items = [
    "FastAPI · DB · OpenAI 연동",
    "MCP: filesystem  /  postgres  /  github",
    "  └ postgres: DB 스키마 직접 조회·검증",
    "Skills: claude-api  /  simplify",
]
for i, item in enumerate(be_items):
    add_text(slide, item,
             Inches(0.6), Inches(3.98 + i * 0.5), Inches(5.2), Inches(0.45),
             font_size=13, color=WHITE if i < 2 else SLATE)

# Frontend Agent 박스
add_rect(slide, Inches(7.3), Inches(3.4), Inches(5.6), Inches(2.8), rgb(0x1E, 0x29, 0x3B))
add_rect(slide, Inches(7.3), Inches(3.4), Inches(5.6), Inches(0.5), rgb(0x10, 0xB9, 0x81))
add_text(slide, "Frontend Agent  (Claude Sonnet 4.6)",
         Inches(7.4), Inches(3.43), Inches(5.4), Inches(0.42),
         font_size=15, bold=True, color=DARK, align=PP_ALIGN.CENTER)

fe_items = [
    "React UI · 컴포넌트 · API 연동",
    "MCP: filesystem  /  github  /  puppeteer",
    "  └ puppeteer: UI 자동 스크린샷 검증",
    "Skills: simplify",
]
for i, item in enumerate(fe_items):
    add_text(slide, item,
             Inches(7.5), Inches(3.98 + i * 0.5), Inches(5.2), Inches(0.45),
             font_size=13, color=WHITE if i < 2 else SLATE)

# 하단 핵심 포인트
add_rect(slide, Inches(0.4), Inches(6.4), Inches(12.5), Inches(0.06), ACCENT)
add_text(slide,
         "worktree 격리로 충돌 없는 병렬 작업  ·  모든 작업이 git commit으로 자동 추적  ·  PM Agent가 PR 자동 리뷰",
         Inches(0.4), Inches(6.55), Inches(12.5), Inches(0.5),
         font_size=14, color=ACCENT, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════
# SLIDE 8 (구 7) — 시스템 아키텍처
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, Inches(1.4), PNU_BLUE)
add_text(slide, "시스템 아키텍처", Inches(0.6), Inches(0.35),
         Inches(11), Inches(0.7), font_size=36, bold=True, color=WHITE)
add_text(slide, "SYSTEM ARCHITECTURE", Inches(0.6), Inches(0.95),
         Inches(11), Inches(0.4), font_size=16, color=ACCENT)

add_rect(slide, 0, Inches(1.4), W, H - Inches(1.4), DARK)

layers = [
    ("Frontend", "React 19 + Vite + TypeScript", Inches(0.5)),
    ("Backend API", "FastAPI (Python) + JWT Auth", Inches(3.2)),
    ("Database", "PostgreSQL + SQLAlchemy", Inches(5.9)),
    ("AI Layer", "OpenAI API (GPT-4o)", Inches(8.6)),
]
for title, sub, y in layers:
    add_rect(slide, Inches(0.4), y, Inches(12.5), Inches(0.06), ACCENT)
    add_text(slide, title, Inches(0.5), y + Inches(0.1),
             Inches(4), Inches(0.4), font_size=18, bold=True, color=ACCENT)
    add_text(slide, sub, Inches(4.5), y + Inches(0.1),
             Inches(8), Inches(0.4), font_size=18, color=LIGHT_GRAY)

add_text(slide, "배포: Vercel (FE)  +  VPS / Docker (BE)",
         Inches(0.5), Inches(6.3), Inches(12), Inches(0.5),
         font_size=18, color=SLATE, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════
# SLIDE 8 — MVP 범위 & 개발 일정
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, Inches(1.4), PNU_BLUE)
add_text(slide, "MVP 범위 & 개발 일정", Inches(0.6), Inches(0.35),
         Inches(11), Inches(0.7), font_size=36, bold=True, color=WHITE)
add_text(slide, "ROADMAP", Inches(0.6), Inches(0.95),
         Inches(11), Inches(0.4), font_size=16, color=ACCENT)

add_rect(slide, 0, Inches(1.4), W, H - Inches(1.4), LIGHT_GRAY)

mvp_items = [
    "✅ 학생 회원가입 / 로그인 (완성)",
    "✅ 학과별 졸업요건 DB + 관리자 CRUD (완성)",
    "🔄 수강 이력 입력 및 조회",
    "🔄 졸업요건 충족률 계산 대시보드",
    "🔄 AI 챗봇 (OpenAI 연동)",
    "🔄 다음 학기 수강 추천 로직",
]
add_text(slide, "우선 구현 기능 (MVP)", Inches(0.5), Inches(1.5),
         Inches(6), Inches(0.5), font_size=20, bold=True, color=PNU_BLUE)
add_bullet_box(slide, mvp_items, Inches(0.5), Inches(2.0),
               Inches(6.0), Inches(4.5), font_size=16,
               color=DARK, bg=WHITE, bullet="")

milestones = [
    ("5월 4~11일", "예선 서류 제출"),
    ("5월 19일", "예선 결과 발표"),
    ("6월 1~4주", "수강이력 + AI 연동 구현"),
    ("7월 1~2주", "What-if + 이력서 기능"),
    ("7월 3~4주", "중간보고회 준비"),
    ("8월 1~4주", "최종 마무리 + 배포"),
    ("8월 28일", "최종 발표 (농심호텔)"),
]
add_text(slide, "개발 일정", Inches(7.0), Inches(1.5),
         Inches(6), Inches(0.5), font_size=20, bold=True, color=PNU_BLUE)
for i, (date, task) in enumerate(milestones):
    y = Inches(2.0 + i * 0.65)
    add_rect(slide, Inches(7.0), y, Inches(2.2), Inches(0.5), PNU_BLUE)
    add_text(slide, date, Inches(7.0), y, Inches(2.2), Inches(0.5),
             font_size=13, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, task, Inches(9.35), y, Inches(3.8), Inches(0.5),
             font_size=14, color=DARK)


# ════════════════════════════════════════════════════════
# SLIDE 9 — 기대효과 & 차별성
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, Inches(1.4), PNU_BLUE)
add_text(slide, "기대효과 & 차별성", Inches(0.6), Inches(0.35),
         Inches(11), Inches(0.7), font_size=36, bold=True, color=WHITE)
add_text(slide, "IMPACT & DIFFERENTIATION", Inches(0.6), Inches(0.95),
         Inches(11), Inches(0.4), font_size=16, color=ACCENT)

add_rect(slide, 0, Inches(1.4), W, H - Inches(1.4), DARK)

diff_items = [
    "학업·졸업·진로·이력서를 하나의 플랫폼에서 통합 관리",
    "학생별 이수 현황 기반 맞춤형 수강 로드맵",
    "정형 데이터 분석 + 생성형 AI 설명 결합 (신뢰성 확보)",
    "정적 PDF가 아닌 인터랙티브 로드맵 UI",
    "복수전공·전과 시나리오 분석 지원",
]
add_text(slide, "차별성", Inches(0.5), Inches(1.6),
         Inches(6), Inches(0.5), font_size=22, bold=True, color=ACCENT)
add_bullet_box(slide, diff_items, Inches(0.5), Inches(2.2),
               Inches(6.0), Inches(4.8), font_size=16,
               color=WHITE, bg=None)

effect_items = [
    "수강 계획 수립 시간 단축",
    "졸업요건 누락 방지",
    "진로 준비 체계화",
    "편입·복학생 적응 지원",
    "취업 자료 정리 부담 완화",
]
add_text(slide, "기대효과", Inches(7.0), Inches(1.6),
         Inches(6), Inches(0.5), font_size=22, bold=True, color=ACCENT)
add_bullet_box(slide, effect_items, Inches(7.0), Inches(2.2),
               Inches(6.0), Inches(4.8), font_size=16,
               color=WHITE, bg=None)


# ════════════════════════════════════════════════════════
# SLIDE 10 — 팀 구성
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, Inches(1.4), PNU_BLUE)
add_text(slide, "팀 구성", Inches(0.6), Inches(0.35),
         Inches(11), Inches(0.7), font_size=36, bold=True, color=WHITE)
add_text(slide, "TEAM", Inches(0.6), Inches(0.95),
         Inches(11), Inches(0.4), font_size=16, color=ACCENT)

add_rect(slide, 0, Inches(1.4), W, H - Inches(1.4), LIGHT_GRAY)

add_text(slide,
         "SW전공자 + 비전공자 각 1명 이상 포함  |  팀원: [추가 필요]",
         Inches(0.5), Inches(1.6), Inches(12), Inches(0.5),
         font_size=18, color=SLATE, align=PP_ALIGN.CENTER)

roles = [
    ("PM / 기획", "SW전공자", [
        "프로젝트 관리",
        "문서·발표자료 작성",
        "AI 도구 총괄 활용",
        "GitHub 관리",
    ]),
    ("Backend", "SW전공자", [
        "FastAPI 서버 개발",
        "DB 설계 및 구현",
        "추천 로직 구현",
        "OpenAI API 연동",
    ]),
    ("Frontend", "SW전공자 또는 비전공자", [
        "React UI 개발",
        "화면 설계·구현",
        "API 연동",
        "UX 개선",
    ]),
    ("기획·데이터", "비전공자", [
        "사용자 요구사항 도출",
        "데이터 정리 (교육과정)",
        "시연 기획",
        "발표 보조",
    ]),
]

for i, (role, spec, items) in enumerate(roles):
    x = Inches(0.4 + i * 3.2)
    add_rect(slide, x, Inches(2.2), Inches(3.0), Inches(4.8), WHITE)
    add_rect(slide, x, Inches(2.2), Inches(3.0), Inches(0.6), PNU_BLUE)
    add_text(slide, role, x, Inches(2.25), Inches(3.0), Inches(0.5),
             font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, spec, x, Inches(2.85), Inches(3.0), Inches(0.35),
             font_size=12, color=ACCENT, align=PP_ALIGN.CENTER)
    for j, item in enumerate(items):
        add_text(slide, f"• {item}", x + Inches(0.15),
                 Inches(3.3 + j * 0.55), Inches(2.7), Inches(0.5),
                 font_size=14, color=DARK)


# ════════════════════════════════════════════════════════
# SLIDE 11 — 마무리
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, H, DARK)
add_rect(slide, 0, Inches(3.3), W, Inches(0.06), ACCENT)

add_text(slide, "PNU-Pathfinder",
         Inches(1), Inches(1.5), Inches(11), Inches(1.2),
         font_size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(slide,
         '"지금 내 상황에서 무엇을 해야 하는가"\n를 AI가 구체적으로 제안합니다.',
         Inches(1), Inches(3.5), Inches(11), Inches(1.4),
         font_size=24, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

add_text(slide, "감사합니다",
         Inches(1), Inches(5.3), Inches(11), Inches(0.7),
         font_size=28, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


# ── 저장 ────────────────────────────────────────────────
out_dir = os.path.join(os.path.dirname(__file__), "..", "output")
os.makedirs(out_dir, exist_ok=True)
out_path = os.path.join(out_dir, "발표자료_PNU-Pathfinder.pptx")
prs.save(out_path)
print(f"✅ PPT 생성 완료: {out_path}")
